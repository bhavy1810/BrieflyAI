# =========================================================
# BrieflyAI – RAG-based PDF, DOCX, TXT, PPT, CSV, Excel, YouTube, Video, Audio & Webpage Summarizer
# LangChain + Groq + FAISS
# UPDATED: Enhanced error handling and dependency checks
# =========================================================

import streamlit as st
import os
import io
import re
import subprocess
import tempfile
from io import BytesIO
from dotenv import load_dotenv
import math
import pandas as pd
import requests
from bs4 import BeautifulSoup
from urllib.parse import urlparse

# PDF
try:
    import PyPDF2
except ImportError:
    import pypdf as PyPDF2

# DOCX
from docx import Document

# PPT/PPTX
from pptx import Presentation

# PDF Export
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet

# LangChain
from langchain_groq import ChatGroq
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS

# YouTube & Transcription - with error handling
try:
    from youtube_transcript_api import YouTubeTranscriptApi
    YOUTUBE_TRANSCRIPT_AVAILABLE = True
except ImportError:
    YOUTUBE_TRANSCRIPT_AVAILABLE = False

try:
    from groq import Groq
    GROQ_AVAILABLE = True
except ImportError:
    GROQ_AVAILABLE = False

# =========================================================
# CONFIG
# =========================================================
load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY")

def load_css(file_name):
    try:
        with open(file_name) as f:
            st.markdown(f"<style>{f.read()}</style>", unsafe_allow_html=True)
    except:
        pass
        
load_css("style.css")
bubbles_html = ''.join([
    f'<div class="floating-bubble bubble-{i}"></div>' 
    for i in range(1, 21)
])
st.markdown(bubbles_html, unsafe_allow_html=True)


st.set_page_config(
    page_title="BrieflyAI",
    page_icon="📄",
    layout="wide"
)

st.title("📄 BrieflyAI")
st.markdown("Upload **PDF / DOCX / TXT / PPT / CSV / Excel / Videos / Audio**, paste text, summarize **YouTube videos** or **webpages** with **RAG-powered AI summaries** 🚀")

# =========================================================
# DEPENDENCY CHECKS
# =========================================================
def check_dependencies():
    """Check if all required dependencies are installed"""
    missing_deps = []
    
    if not YOUTUBE_TRANSCRIPT_AVAILABLE:
        missing_deps.append(("youtube-transcript-api", "pip install youtube-transcript-api --break-system-packages"))
    
    if not GROQ_AVAILABLE:
        missing_deps.append(("groq", "pip install groq --break-system-packages"))
    
    # Check yt-dlp
    try:
        subprocess.run(["yt-dlp", "--version"], capture_output=True, check=True)
    except (FileNotFoundError, subprocess.CalledProcessError):
        missing_deps.append(("yt-dlp", "pip install yt-dlp --break-system-packages"))
    
    # Check ffmpeg
    try:
        subprocess.run(["ffmpeg", "-version"], capture_output=True, check=True)
    except (FileNotFoundError, subprocess.CalledProcessError):
        missing_deps.append(("ffmpeg", "sudo apt-get install ffmpeg"))
    
    return missing_deps

def show_dependency_warning():
    """Show warning about missing dependencies for YouTube features"""
    missing = check_dependencies()
    
    if missing and (st.session_state.get("show_youtube_warning", True)):
        with st.expander("⚠️ YouTube Features - Missing Dependencies", expanded=True):
            st.warning("Some YouTube features may not work. Missing dependencies:")
            for dep, install_cmd in missing:
                st.code(f"❌ {dep}\n   Install: {install_cmd}")
            
            if not GROQ_API_KEY:
                st.error("🔑 GROQ_API_KEY not set!")
                st.info("Get your free API key from: https://console.groq.com/keys")
                st.code("export GROQ_API_KEY='your_key_here'\n# OR add to .env file:\nGROQ_API_KEY=your_key_here")
            
            if st.button("✅ I've installed the dependencies"):
                st.session_state.show_youtube_warning = False
                st.rerun()

# =========================================================
# SESSION STATE
# =========================================================
defaults = {
    "quick_summary": None,
    "detailed_summary": None,
    "vectorstore": None,
    "summary_type": None,
    "is_generating": False,
    "source_text": None,
    "show_youtube_warning": True,
}

for k, v in defaults.items():
    if k not in st.session_state:
        st.session_state[k] = v

# =========================================================
# LLM SETUP
# =========================================================
if not GROQ_API_KEY:
    st.error("❌ GROQ_API_KEY not found in environment variables.")
    st.info("📝 To fix this:")
    st.code("1. Get API key from: https://console.groq.com/keys\n2. Add to .env file: GROQ_API_KEY=your_key\n3. Or run: export GROQ_API_KEY='your_key'")
    st.stop()

llm = ChatGroq(
    api_key=GROQ_API_KEY,
    model="llama-3.3-70b-versatile",
    temperature=0.3,
)
output_parser = StrOutputParser()

# Setup Groq for transcription
if GROQ_AVAILABLE and GROQ_API_KEY:
    groq_client = Groq(api_key=GROQ_API_KEY)
else:
    groq_client = None

# =========================================================
# HELPERS - FILE PROCESSING
# =========================================================
def extract_text_from_pdf(file):
    reader = PyPDF2.PdfReader(io.BytesIO(file.read()))
    return "\n".join(p.extract_text() or "" for p in reader.pages)

def extract_text_from_docx(file):
    doc = Document(file)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def extract_text_from_txt(file):
    """Extract text from TXT file with encoding detection."""
    try:
        content = file.read()
        # Try UTF-8 first
        try:
            return content.decode('utf-8')
        except UnicodeDecodeError:
            # Fallback to latin-1
            return content.decode('latin-1')
    except Exception as e:
        st.error(f"Error reading TXT file: {str(e)}")
        return ""

def extract_text_from_ppt(file):
    """Extract text from PPT/PPTX file."""
    try:
        prs = Presentation(file)
        text_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # Extract text from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip():
                            slide_text.append(row_text)
            
            text_content.append("\n".join(slide_text))
        
        return "\n\n".join(text_content)
    except Exception as e:
        st.error(f"Error reading PPT file: {str(e)}")
        return ""

def extract_text_from_csv(file):
    """Extract and format text from CSV file."""
    try:
        # Try to read CSV with different encodings
        try:
            df = pd.read_csv(file, encoding='utf-8')
        except UnicodeDecodeError:
            file.seek(0)
            df = pd.read_csv(file, encoding='latin-1')
        
        # Create formatted text
        text_content = []
        text_content.append(f"CSV Data Summary:")
        text_content.append(f"Total Rows: {len(df)}")
        text_content.append(f"Total Columns: {len(df.columns)}")
        text_content.append(f"\nColumns: {', '.join(df.columns.tolist())}\n")
        text_content.append("=" * 50)
        text_content.append("\nData Content:\n")
        
        # Convert dataframe to formatted text
        for idx, row in df.iterrows():
            row_text = f"Row {idx + 1}:\n"
            for col in df.columns:
                value = row[col]
                if pd.notna(value):
                    row_text += f"  {col}: {value}\n"
            text_content.append(row_text)
        
        return "\n".join(text_content)
    except Exception as e:
        st.error(f"Error reading CSV file: {str(e)}")
        return ""

def extract_text_from_excel(file):
    """Extract and format text from Excel file."""
    try:
        # Read all sheets
        excel_file = pd.ExcelFile(file)
        text_content = []
        
        text_content.append(f"Excel File Summary:")
        text_content.append(f"Total Sheets: {len(excel_file.sheet_names)}")
        text_content.append(f"Sheet Names: {', '.join(excel_file.sheet_names)}\n")
        text_content.append("=" * 50)
        
        # Process each sheet
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file, sheet_name=sheet_name)
            
            text_content.append(f"\n\n--- Sheet: {sheet_name} ---")
            text_content.append(f"Rows: {len(df)} | Columns: {len(df.columns)}")
            text_content.append(f"Columns: {', '.join(df.columns.tolist())}\n")
            
            # Convert dataframe to formatted text
            for idx, row in df.iterrows():
                row_text = f"Row {idx + 1}:\n"
                for col in df.columns:
                    value = row[col]
                    if pd.notna(value):
                        row_text += f"  {col}: {value}\n"
                text_content.append(row_text)
            
            # Reset file pointer for next sheet
            file.seek(0)
        
        return "\n".join(text_content)
    except Exception as e:
        st.error(f"Error reading Excel file: {str(e)}")
        return ""

# =========================================================
# HELPERS - WEBPAGE SCRAPING
# =========================================================
def is_valid_url(url):
    """Validate URL format."""
    try:
        result = urlparse(url)
        return all([result.scheme, result.netloc])
    except:
        return False

def extract_text_from_webpage(url):
    """
    Extract main text content from a webpage.
    Removes scripts, styles, navigation, and other non-content elements.
    """
    try:
        # Set headers to mimic a browser
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        
        # Fetch the webpage
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        
        # Parse with BeautifulSoup
        soup = BeautifulSoup(response.content, 'html.parser')
        
        # Remove unwanted elements
        for element in soup(['script', 'style', 'nav', 'header', 'footer', 'aside', 'iframe', 'noscript']):
            element.decompose()
        
        # Extract title
        title = soup.find('title')
        title_text = title.get_text().strip() if title else "No Title"
        
        # Extract meta description
        meta_desc = soup.find('meta', attrs={'name': 'description'})
        description = meta_desc.get('content', '').strip() if meta_desc else ""
        
        # Extract main content
        # Try to find main content area
        main_content = soup.find('main') or soup.find('article') or soup.find('div', class_=re.compile('content|main|article'))
        
        if main_content:
            text_elements = main_content.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'li'])
        else:
            text_elements = soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'li'])
        
        # Extract text from elements
        content_parts = []
        
        # Add title and description
        content_parts.append(f"TITLE: {title_text}\n")
        if description:
            content_parts.append(f"DESCRIPTION: {description}\n")
        content_parts.append(f"URL: {url}\n")
        content_parts.append("=" * 50 + "\n")
        
        # Add main content
        for element in text_elements:
            text = element.get_text().strip()
            if text and len(text) > 20:  # Filter out very short text
                # Add heading markers
                if element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                    content_parts.append(f"\n## {text}\n")
                else:
                    content_parts.append(text)
        
        full_text = "\n".join(content_parts)
        
        # Clean up extra whitespace
        full_text = re.sub(r'\n{3,}', '\n\n', full_text)
        
        return full_text
        
    except requests.exceptions.Timeout:
        raise Exception("Request timed out. The webpage took too long to respond.")
    except requests.exceptions.ConnectionError:
        raise Exception("Connection error. Could not connect to the webpage.")
    except requests.exceptions.HTTPError as e:
        raise Exception(f"HTTP error: {e.response.status_code}")
    except Exception as e:
        raise Exception(f"Failed to extract content: {str(e)}")

# =========================================================
# HELPERS - VECTORSTORE & SUMMARIZATION
# =========================================================
def build_vectorstore(text):
    splitter = RecursiveCharacterTextSplitter(chunk_size=1200, chunk_overlap=200)
    docs = splitter.create_documents([text])
    embeddings = HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2"
    )
    return FAISS.from_documents(docs, embeddings)

def rag_summary(vectorstore, summary_type):
    if summary_type == "quick":
        k = 4
        max_words = 200
        template = """
You are an expert summarization AI.
Using ONLY the context below, generate a clear and concise summary of about {max_words} words.

Context:
{context}

Quick Summary:
"""
    else:
        k = 10
        max_words = 800
        template = """
You are an expert summarization AI.
Using ONLY the context below, generate a detailed summary of about {max_words} words.
Use headings, bullet points, and clear explanations.

Context:
{context}

Detailed Summary:
"""

    retriever = vectorstore.as_retriever(search_kwargs={"k": k})
    docs = retriever.invoke("Summarize the document")
    context = "\n\n".join(d.page_content for d in docs)

    prompt = ChatPromptTemplate.from_template(template)
    chain = prompt | llm | output_parser

    return chain.invoke({
        "context": context[:25000],
        "max_words": max_words
    })

# =========================================================
# HELPERS - YOUTUBE
# =========================================================
def extract_video_id(url):
    patterns = [
        r"v=([^&]+)",
        r"youtu\.be/([^?&]+)",
        r"shorts/([^?&]+)",
        r"embed/([^?&]+)"
    ]
    for pattern in patterns:
        match = re.search(pattern, url)
        if match:
            return match.group(1)
    return None

def get_youtube_transcript(video_id):
    """Get transcript from YouTube if available"""
    if not YOUTUBE_TRANSCRIPT_AVAILABLE:
        return None
    
    try:
        transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
        try:
            transcript = transcript_list.find_manually_created_transcript(['en'])
        except:
            try:
                transcript = transcript_list.find_generated_transcript(['en'])
            except:
                transcript = next(iter(transcript_list._transcripts.values()))
        
        data = transcript.fetch()
        return " ".join(item["text"] for item in data)
    except Exception as e:
        return None

def get_audio_duration(audio_path):
    """Get audio duration using ffprobe."""
    try:
        result = subprocess.run(
            [
                "ffprobe",
                "-v", "error",
                "-show_entries", "format=duration",
                "-of", "default=noprint_wrappers=1:nokey=1",
                audio_path
            ],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True
        )
        return float(result.stdout.strip())
    except:
        return None

def split_audio_with_ffmpeg(audio_path, max_size_mb=20):
    """
    Split audio file into chunks using FFmpeg (no pydub needed).
    Returns list of file paths.
    """
    file_size_mb = os.path.getsize(audio_path) / (1024 * 1024)
    
    if file_size_mb <= max_size_mb:
        return [audio_path]
    
    # Get audio duration
    duration = get_audio_duration(audio_path)
    if not duration:
        st.warning("⚠️ Could not determine audio duration. Using single file.")
        return [audio_path]
    
    # Calculate number of chunks needed
    num_chunks = math.ceil(file_size_mb / max_size_mb)
    chunk_duration = duration / num_chunks
    
    chunk_paths = []
    base_path = audio_path.replace(".mp3", "")
    
    for i in range(num_chunks):
        start_time = i * chunk_duration
        chunk_path = f"{base_path}_chunk_{i}.mp3"
        
        # Use FFmpeg to extract chunk
        subprocess.run(
            [
                "ffmpeg",
                "-i", audio_path,
                "-ss", str(start_time),
                "-t", str(chunk_duration),
                "-acodec", "copy",
                "-y",  # Overwrite output file
                chunk_path
            ],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL
        )
        
        if os.path.exists(chunk_path):
            chunk_paths.append(chunk_path)
    
    return chunk_paths if chunk_paths else [audio_path]

def transcribe_audio_chunks(audio_path):
    """
    Transcribe audio file with automatic chunking if needed.
    """
    if not groq_client:
        raise Exception("Groq client not available. Please install groq and set GROQ_API_KEY.")
    
    # Check file size and split if necessary
    file_size_mb = os.path.getsize(audio_path) / (1024 * 1024)
    
    if file_size_mb > 20:
        st.info(f"📦 Large audio file ({file_size_mb:.1f}MB) - splitting into chunks...")
        audio_chunks = split_audio_with_ffmpeg(audio_path, max_size_mb=20)
    else:
        audio_chunks = [audio_path]
    
    # Transcribe each chunk
    full_transcription = []
    
    for idx, chunk_path in enumerate(audio_chunks):
        if len(audio_chunks) > 1:
            chunk_size = os.path.getsize(chunk_path) / (1024 * 1024)
            st.info(f"🎙️ Transcribing chunk {idx + 1}/{len(audio_chunks)} ({chunk_size:.1f}MB)...")
        
        try:
            with open(chunk_path, "rb") as audio_file:
                transcription = groq_client.audio.transcriptions.create(
                    file=(os.path.basename(chunk_path), audio_file.read()),
                    model="whisper-large-v3-turbo",
                    response_format="text"
                )
                full_transcription.append(transcription)
        except Exception as e:
            st.error(f"❌ Error transcribing chunk {idx + 1}: {str(e)}")
            continue
    
    if not full_transcription:
        raise Exception("All transcription attempts failed")
    
    return " ".join(full_transcription)

def groq_whisper_transcription(video_url):
    """
    Enhanced transcription with FFmpeg-based audio splitting for YouTube videos.
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        audio_path = os.path.join(tmpdir, "audio.mp3")
        
        # Download audio with size optimization
        result = subprocess.run(
            [
                "yt-dlp",
                "-f", "bestaudio[ext=m4a]/bestaudio",
                "-x",
                "--audio-format", "mp3",
                "--audio-quality", "5",  # Compress audio
                "-o", audio_path,
                video_url
            ],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL
        )
        
        if result.returncode != 0:
            raise Exception("Failed to download audio from YouTube")
        
        return transcribe_audio_chunks(audio_path)

def transcribe_local_video(video_file):
    """
    Transcribe local video file using Groq Whisper with chunking support.
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        # Save uploaded video
        video_path = os.path.join(tmpdir, f"video{os.path.splitext(video_file.name)[1]}")
        with open(video_path, "wb") as f:
            f.write(video_file.read())
        
        # Extract audio from video
        audio_path = os.path.join(tmpdir, "audio.mp3")
        
        result = subprocess.run(
            [
                "ffmpeg",
                "-i", video_path,
                "-vn",  # No video
                "-acodec", "libmp3lame",
                "-q:a", "5",  # Compress audio
                "-y",
                audio_path
            ],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL
        )
        
        if result.returncode != 0:
            raise Exception("Failed to extract audio from video")
        
        return transcribe_audio_chunks(audio_path)

def transcribe_audio_file(audio_file):
    """
    Transcribe audio file using Groq Whisper with chunking support.
    Converts to MP3 if needed.
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        # Save uploaded audio
        file_ext = os.path.splitext(audio_file.name)[1].lower()
        input_path = os.path.join(tmpdir, f"input{file_ext}")
        
        with open(input_path, "wb") as f:
            f.write(audio_file.read())
        
        # Convert to MP3 if not already MP3
        if file_ext != ".mp3":
            audio_path = os.path.join(tmpdir, "audio.mp3")
            result = subprocess.run(
                [
                    "ffmpeg",
                    "-i", input_path,
                    "-acodec", "libmp3lame",
                    "-q:a", "5",  # Compress audio
                    "-y",
                    audio_path
                ],
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL
            )
            
            if result.returncode != 0:
                raise Exception("Failed to convert audio file")
        else:
            audio_path = input_path
        
        return transcribe_audio_chunks(audio_path)

# =========================================================
# EXPORTS
# =========================================================
def export_pdf(text):
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4)
    styles = getSampleStyleSheet()
    story = []

    for line in text.split("\n"):
        story.append(Paragraph(line, styles["Normal"]))
        story.append(Spacer(1, 6))

    doc.build(story)
    buffer.seek(0)
    return buffer

def export_docx(text):
    buffer = BytesIO()
    doc = Document()
    for line in text.split("\n"):
        doc.add_paragraph(line)
    doc.save(buffer)
    buffer.seek(0)
    return buffer

# =========================================================
# SHOW DEPENDENCY WARNING
# =========================================================
show_dependency_warning()

# =========================================================
# UI - TABS
# =========================================================
tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs([
    "📄 Upload File", 
    "✍️ Paste Text", 
    "🌐 Webpage", 
    "📺 YouTube Video", 
    "🎬 Local Video", 
    "🎵 Audio File"
])

# =========================================================
# TAB 1: FILE UPLOAD
# =========================================================
with tab1:
    st.markdown("### 📄 Upload Document")
    st.info("📄 Summarize documents by uploading them below.")
    
    uploaded = st.file_uploader(
        "Choose a file",
        type=["pdf", "docx", "txt", "ppt", "pptx", "csv", "xls", "xlsx"]
    )
    
    col1, col2 = st.columns(2)

    with col1:
        if st.button("⚡ Quick Summary", key="btn_quick_file", use_container_width=True):
            if not uploaded:
                st.warning("Please upload a file")
            else:
                file_ext = uploaded.name.split(".")[-1].lower()
                
                # Display file info
                file_size_mb = len(uploaded.getvalue()) / (1024 * 1024)
                st.caption(f"📦 File: {uploaded.name} ({file_size_mb:.2f} MB)")
                
                # Extract text based on file type
                with st.spinner(f"📖 Reading {file_ext.upper()} file..."):
                    if file_ext == "pdf":
                        text = extract_text_from_pdf(uploaded)
                    elif file_ext == "docx":
                        text = extract_text_from_docx(uploaded)
                    elif file_ext == "txt":
                        text = extract_text_from_txt(uploaded)
                    elif file_ext in ["ppt", "pptx"]:
                        text = extract_text_from_ppt(uploaded)
                    elif file_ext == "csv":
                        text = extract_text_from_csv(uploaded)
                    elif file_ext in ["xls", "xlsx"]:
                        text = extract_text_from_excel(uploaded)
                    else:
                        st.error("Unsupported file type")
                        st.stop()
                
                if text:
                    word_count = len(text.split())
                    st.success(f"✅ Extracted {word_count:,} words")
                    
                    if not st.session_state.is_generating:
                        st.session_state.is_generating = True
                        st.session_state.source_text = text
                        with st.spinner("Indexing & Summarizing..."):
                            st.session_state.vectorstore = build_vectorstore(text)
                            st.session_state.quick_summary = rag_summary(
                                st.session_state.vectorstore, "quick"
                            )
                        st.session_state.summary_type = "quick"
                        st.session_state.is_generating = False
                        st.rerun()
                else:
                    st.warning("⚠️ No text extracted from file")
                    st.stop()

    with col2:
        if st.button("📋 Detailed Summary", key="btn_detailed_file", use_container_width=True):
            if not uploaded:
                st.warning("Please upload a file")
            else:
                file_ext = uploaded.name.split(".")[-1].lower()
                
                # Display file info
                file_size_mb = len(uploaded.getvalue()) / (1024 * 1024)
                st.caption(f"📦 File: {uploaded.name} ({file_size_mb:.2f} MB)")
                
                # Extract text based on file type
                with st.spinner(f"📖 Reading {file_ext.upper()} file..."):
                    if file_ext == "pdf":
                        text = extract_text_from_pdf(uploaded)
                    elif file_ext == "docx":
                        text = extract_text_from_docx(uploaded)
                    elif file_ext == "txt":
                        text = extract_text_from_txt(uploaded)
                    elif file_ext in ["ppt", "pptx"]:
                        text = extract_text_from_ppt(uploaded)
                    elif file_ext == "csv":
                        text = extract_text_from_csv(uploaded)
                    elif file_ext in ["xls", "xlsx"]:
                        text = extract_text_from_excel(uploaded)
                    else:
                        st.error("Unsupported file type")
                        st.stop()
                
                if text:
                    word_count = len(text.split())
                    st.success(f"✅ Extracted {word_count:,} words")
                    
                    if not st.session_state.is_generating:
                        st.session_state.is_generating = True
                        st.session_state.source_text = text
                        with st.spinner("Indexing & Summarizing..."):
                            st.session_state.vectorstore = build_vectorstore(text)
                            st.session_state.detailed_summary = rag_summary(
                                st.session_state.vectorstore, "detailed"
                            )
                        st.session_state.summary_type = "detailed"
                        st.session_state.is_generating = False
                        st.rerun()
                else:
                    st.warning("⚠️ No text extracted from file")
                    st.stop()

# =========================================================
# TAB 2: PASTE TEXT
# =========================================================
with tab2:
    st.markdown("### ✍️ Summarize Text")
    st.info("✍️ Paste or type any text to summarize.")
    pasted = st.text_area("✍️ Paste text", height=200)
    col1, col2 = st.columns(2)

    with col1:
        if st.button("⚡ Quick Summary", key="btn_quick_paste", use_container_width=True):
            if not pasted:
                st.warning("Please paste some text")
            else:
                if not st.session_state.is_generating:
                    st.session_state.is_generating = True
                    st.session_state.source_text = pasted
                    with st.spinner("Indexing & Summarizing..."):
                        st.session_state.vectorstore = build_vectorstore(pasted)
                        st.session_state.quick_summary = rag_summary(
                            st.session_state.vectorstore, "quick"
                        )
                    st.session_state.summary_type = "quick"
                    st.session_state.is_generating = False
                    st.rerun()

    with col2:
        if st.button("📋 Detailed Summary", key="btn_detailed_paste", use_container_width=True):
            if not pasted:
                st.warning("Please paste some text")
            else:
                if not st.session_state.is_generating:
                    st.session_state.is_generating = True
                    st.session_state.source_text = pasted
                    with st.spinner("Indexing & Summarizing..."):
                        st.session_state.vectorstore = build_vectorstore(pasted)
                        st.session_state.detailed_summary = rag_summary(
                            st.session_state.vectorstore, "detailed"
                        )
                    st.session_state.summary_type = "detailed"
                    st.session_state.is_generating = False
                    st.rerun()
                    
# =========================================================
# TAB 3: WEBPAGE SUMMARIZER
# =========================================================
with tab3:
    st.markdown("### 🌐 Summarize Webpage")
    st.info("🔗 Enter any webpage URL to extract and summarize its content.")
    
    webpage_url = st.text_input(
        "🔗 Enter webpage URL",
        key="webpage_url_input",
        placeholder="https://example.com/article"
    )
    
    # Show URL validation
    if webpage_url:
        if is_valid_url(webpage_url):
            st.success("✅ Valid URL")
        else:
            st.error("❌ Invalid URL format. Please include http:// or https://")
    
    col1, col2 = st.columns(2)
    
    with col1:
        if st.button("⚡ Quick Summary", key="btn_webpage_quick", use_container_width=True):
            if not webpage_url:
                st.warning("Please enter a webpage URL")
            elif not is_valid_url(webpage_url):
                st.error("Invalid URL format")
            else:
                with st.spinner("🌐 Fetching and analyzing webpage..."):
                    try:
                        # Extract webpage content
                        webpage_text = extract_text_from_webpage(webpage_url)
                        
                        if not webpage_text or len(webpage_text.strip()) < 100:
                            st.error("⚠️ Could not extract sufficient content from the webpage")
                            st.stop()
                        
                        word_count = len(webpage_text.split())
                        st.success(f"✅ Extracted {word_count:,} words from webpage")
                        
                        # Build vectorstore and generate summary
                        st.session_state.source_text = webpage_text
                        st.session_state.vectorstore = build_vectorstore(webpage_text)
                        st.session_state.quick_summary = rag_summary(
                            st.session_state.vectorstore, "quick"
                        )
                        st.session_state.summary_type = "quick"
                        st.rerun()
                        
                    except Exception as e:
                        st.error(f"❌ Error: {str(e)}")
                        st.info("💡 Tip: Make sure the URL is accessible and contains readable content")
    
    with col2:
        if st.button("📋 Detailed Summary", key="btn_webpage_detailed", use_container_width=True):
            if not webpage_url:
                st.warning("Please enter a webpage URL")
            elif not is_valid_url(webpage_url):
                st.error("Invalid URL format")
            else:
                with st.spinner("🌐 Fetching and analyzing webpage..."):
                    try:
                        # Extract webpage content
                        webpage_text = extract_text_from_webpage(webpage_url)
                        
                        if not webpage_text or len(webpage_text.strip()) < 100:
                            st.error("⚠️ Could not extract sufficient content from the webpage")
                            st.stop()
                        
                        word_count = len(webpage_text.split())
                        st.success(f"✅ Extracted {word_count:,} words from webpage")
                        
                        # Build vectorstore and generate summary
                        st.session_state.source_text = webpage_text
                        st.session_state.vectorstore = build_vectorstore(webpage_text)
                        st.session_state.detailed_summary = rag_summary(
                            st.session_state.vectorstore, "detailed"
                        )
                        st.session_state.summary_type = "detailed"
                        st.rerun()
                        
                    except Exception as e:
                        st.error(f"❌ Error: {str(e)}")
                        st.info("💡 Tip: Make sure the URL is accessible and contains readable content")

# =========================================================
# TAB 4: YOUTUBE VIDEO
# =========================================================
with tab4:
    st.markdown("### 📺 Summarize YouTube Video")
    st.info("🔗 Enter any YouTube link to extract and summarize its content.")
    
    # Check if YouTube features are available
    missing = check_dependencies()
    youtube_enabled = not missing and GROQ_API_KEY
    
    if not youtube_enabled:
        st.error("⚠️ YouTube features require additional setup. See warning at the top of the page.")
    
    youtube_link = st.text_input(
        "🔗 Enter YouTube Video Link",
        placeholder="https://www.youtube.com/watch?v=...",
        disabled=not youtube_enabled
    )
    
    video_id = extract_video_id(youtube_link) if youtube_link else None
    
    if video_id:
        st.image(
            f"https://img.youtube.com/vi/{video_id}/0.jpg",
            width=720
        )
    elif youtube_link:
        st.warning("⚠️ Invalid YouTube link")
    
    col1, col2 = st.columns(2)
    
    with col1:
        if st.button("⚡ Quick Summary", key="btn_youtube_quick", use_container_width=True, disabled=not youtube_enabled):
            if not youtube_link:
                st.warning("Please enter a YouTube link")
            elif not video_id:
                st.error("Invalid YouTube URL")
            else:
                with st.spinner("Processing video..."):
                    try:
                        transcript = get_youtube_transcript(video_id)
                        
                        if transcript:
                            st.info("✅ Using YouTube captions")
                        else:
                            try:
                                transcript = groq_whisper_transcription(youtube_link)
                            except Exception as e:
                                st.error(f"❌ Transcription failed: {str(e)}")
                                st.info("💡 Make sure yt-dlp and ffmpeg are installed")
                                st.stop()
                        
                        st.session_state.source_text = transcript
                        st.session_state.vectorstore = build_vectorstore(transcript)
                        st.session_state.quick_summary = rag_summary(
                            st.session_state.vectorstore, "quick"
                        )
                        st.session_state.summary_type = "quick"
                        st.rerun()
                    except Exception as e:
                        st.error(f"❌ Error: {str(e)}")
    
    with col2:
        if st.button("📋 Detailed Summary", key="btn_youtube_detailed", use_container_width=True, disabled=not youtube_enabled):
            if not youtube_link:
                st.warning("Please enter a YouTube link")
            elif not video_id:
                st.error("Invalid YouTube URL")
            else:
                with st.spinner("Processing video with RAG..."):
                    try:
                        transcript = get_youtube_transcript(video_id)
                        
                        if transcript:
                            st.info("✅ Using YouTube captions")
                        else:
                            st.info("🎙️ No captions found — using Groq Whisper")
                            try:
                                transcript = groq_whisper_transcription(youtube_link)
                            except Exception as e:
                                st.error(f"❌ Transcription failed: {str(e)}")
                                st.info("💡 Make sure yt-dlp and ffmpeg are installed")
                                st.stop()
                        
                        st.session_state.source_text = transcript
                        st.session_state.vectorstore = build_vectorstore(transcript)
                        st.session_state.detailed_summary = rag_summary(
                            st.session_state.vectorstore, "detailed"
                        )
                        st.session_state.summary_type = "detailed"
                        st.rerun()
                    except Exception as e:
                        st.error(f"❌ Error: {str(e)}")

# =========================================================
# TAB 5: LOCAL VIDEO
# =========================================================
with tab5:
    st.markdown("### 🎬 Summarize Local Video")
    st.info("🎬 Summarize videos from your device.")
    
    # Check if video features are available
    video_enabled = groq_client is not None
    
    if not video_enabled:
        st.error("⚠️ Video transcription requires Groq API. See warning at the top of the page.")
    
    video_file = st.file_uploader(
        "Choose a video file",
        type=["mp4", "avi", "mov", "mkv", "flv", "wmv"],
        key="local_video_uploader",
        disabled=not video_enabled
    )
    
    if video_file:
        st.video(video_file)
        
        file_size_mb = len(video_file.getvalue()) / (1024 * 1024)
        st.caption(f"📦 File size: {file_size_mb:.1f} MB")
    
    col1, col2 = st.columns(2)
    
    with col1:
        if st.button("⚡ Quick Summary", key="btn_local_video_quick", use_container_width=True, disabled=not video_enabled):
            if not video_file:
                st.warning("Please upload a video file")
            else:
                with st.spinner("🎙️ Transcribing video audio..."):
                    try:
                        # Reset file pointer
                        video_file.seek(0)
                        transcript = transcribe_local_video(video_file)
                        
                        st.success(f"✅ Transcribed {len(transcript.split())} words")
                        
                        st.session_state.source_text = transcript
                        st.session_state.vectorstore = build_vectorstore(transcript)
                        st.session_state.quick_summary = rag_summary(
                            st.session_state.vectorstore, "quick"
                        )
                        st.session_state.summary_type = "quick"
                        st.rerun()
                    except Exception as e:
                        st.error(f"❌ Transcription failed: {str(e)}")
                        st.stop()
    
    with col2:
        if st.button("📋 Detailed Summary", key="btn_local_video_detailed", use_container_width=True, disabled=not video_enabled):
            if not video_file:
                st.warning("Please upload a video file")
            else:
                with st.spinner("🎙️ Transcribing video audio..."):
                    try:
                        # Reset file pointer
                        video_file.seek(0)
                        transcript = transcribe_local_video(video_file)
                        
                        st.success(f"✅ Transcribed {len(transcript.split())} words")
                        
                        st.session_state.source_text = transcript
                        st.session_state.vectorstore = build_vectorstore(transcript)
                        st.session_state.detailed_summary = rag_summary(
                            st.session_state.vectorstore, "detailed"
                        )
                        st.session_state.summary_type = "detailed"
                        st.rerun()
                    except Exception as e:
                        st.error(f"❌ Transcription failed: {str(e)}")
                        st.stop()

# =========================================================
# TAB 6: AUDIO FILE
# =========================================================
with tab6:
    st.markdown("### 🎵 Summarize Audio")
    st.info("🎵 Summarize audio files from your device.")
    
    # Check if audio features are available
    audio_enabled = groq_client is not None
    
    if not audio_enabled:
        st.error("⚠️ Audio transcription requires Groq API. See warning at the top of the page.")
    
    audio_file = st.file_uploader(
        "Choose an audio file",
        type=["mp3", "wav", "m4a", "flac", "ogg", "aac"],
        key="audio_file_uploader",
        disabled=not audio_enabled
    )
    
    if audio_file:
        st.audio(audio_file)
        
        file_size_mb = len(audio_file.getvalue()) / (1024 * 1024)
        st.caption(f"📦 File size: {file_size_mb:.1f} MB")
        
        if file_size_mb > 20:
            st.info("ℹ️ Large audio file will be automatically split into chunks for transcription")
    
    col1, col2 = st.columns(2)
    
    with col1:
        if st.button("⚡ Quick Summary", key="btn_audio_quick", use_container_width=True, disabled=not audio_enabled):
            if not audio_file:
                st.warning("Please upload an audio file")
            else:
                with st.spinner("🎙️ Transcribing audio..."):
                    try:
                        # Reset file pointer
                        audio_file.seek(0)
                        transcript = transcribe_audio_file(audio_file)
                        
                        st.success(f"✅ Transcribed {len(transcript.split())} words")
                        
                        st.session_state.source_text = transcript
                        st.session_state.vectorstore = build_vectorstore(transcript)
                        st.session_state.quick_summary = rag_summary(
                            st.session_state.vectorstore, "quick"
                        )
                        st.session_state.summary_type = "quick"
                        st.rerun()
                    except Exception as e:
                        st.error(f"❌ Transcription failed: {str(e)}")
                        st.stop()
    
    with col2:
        if st.button("📋 Detailed Summary", key="btn_audio_detailed", use_container_width=True, disabled=not audio_enabled):
            if not audio_file:
                st.warning("Please upload an audio file")
            else:
                with st.spinner("🎙️ Transcribing audio..."):
                    try:
                        # Reset file pointer
                        audio_file.seek(0)
                        transcript = transcribe_audio_file(audio_file)
                        
                        st.success(f"✅ Transcribed {len(transcript.split())} words")
                        
                        st.session_state.source_text = transcript
                        st.session_state.vectorstore = build_vectorstore(transcript)
                        st.session_state.detailed_summary = rag_summary(
                            st.session_state.vectorstore, "detailed"
                        )
                        st.session_state.summary_type = "detailed"
                        st.rerun()
                    except Exception as e:
                        st.error(f"❌ Transcription failed: {str(e)}")
                        st.stop()

# =========================================================
# RESULT DISPLAY
# =========================================================
summary = (
    st.session_state.quick_summary
    if st.session_state.summary_type == "quick"
    else st.session_state.detailed_summary
)

if summary:
    st.divider()
    st.header("🧠 Summary")
    
    # Render the Markdown Summary
    st.markdown(summary)

    # Copy Functionality
    st.write("---")
    with st.expander("📝 View Raw Text (Copy to Clipboard)"):
        st.code(summary, language="markdown")
        
    # Exports
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.download_button(
            label="📄 Export TXT",
            data=summary,
            file_name="summary.txt",
            mime="text/plain",
            use_container_width=True
        )
    with col2:
        st.download_button(
            label="📕 Export PDF",
            data=export_pdf(summary),
            file_name="summary.pdf",
            mime="application/pdf",
            use_container_width=True
        )
    with col3:
        st.download_button(
            label="📝 Export DOCX",
            data=export_docx(summary),
            file_name="summary.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            use_container_width=True,
        )

# =========================================================
# FOOTER
# =========================================================
st.divider()